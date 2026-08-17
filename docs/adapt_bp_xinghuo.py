# -*- coding: utf-8 -*-
"""把 business_plan_studymate.pptx 从港科大版适配为星火计划版
- 不改团队配置（5 人名单、不理不器、队长信息均保留）
- Slide 1: 赛事名 → 2026「星火计划」超级个体（OPC）创业挑战赛 · 软件赛道
- 新增 OPC 主题页（插在定价页后、团队页前），页码顺延
- Slide 12(结尾): 三条诉求改为星火权益；赛道名改为软件赛道
"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

SRC = "business_plan_studymate.pptx"           # 港科大原版（保持不动）
DST = "business_plan_studymate_xinghuo.pptx"   # 星火计划版（另存新文件）
C_TEXT = RGBColor(0x11, 0x11, 0x11)
C_MUTED = RGBColor(0x5B, 0x65, 0x73)
C_ACCENT = RGBColor(0x3D, 0x8D, 0xFF)
C_CARD = RGBColor(0xF0, 0xF1, 0xF3)
C_DIVIDER = RGBColor(0xE2, 0xE4, 0xE8)

prs = Presentation(SRC)

# ── 1. 封面赛事名 ────────────────────────────────────────────────
s1 = prs.slides[0]
for sh in s1.shapes:
    if sh.has_text_frame and "香港科技大学" in sh.text_frame.text:
        sh.text_frame.paragraphs[0].runs[0].text = (
            "2026「星火计划」超级个体（OPC）创业挑战赛 · 软件赛道"
        )

# ── 2. 结尾页诉求与赛道 ─────────────────────────────────────────
s_end = prs.slides[11]
replace_map = {
    "01  CPA 种子用户与试点渠道": "01  七天集训与导师打磨，完成产品化冲刺",
    "02  产品化与产业合作指导": "02  投资人闭门对接与融资连接",
    "03  孵化与后续融资连接": "03  松江 OPC 落地扶持与产业资源对接",
    "人工智能赛道 · 初创组": "软件赛道",
}
for sh in s_end.shapes:
    if not sh.has_text_frame:
        continue
    for p in sh.text_frame.paragraphs:
        for r in p.runs:
            if r.text in replace_map:
                r.text = replace_map[r.text]
            elif "人工智能赛道" in r.text:  # 页脚整行兜底
                r.text = r.text.replace("人工智能赛道 · 初创组", "软件赛道")

# ── 3. 新增 OPC 主题页 ──────────────────────────────────────────
layout = prs.slide_masters[0].slide_layouts[0]
slide = prs.slides.add_slide(layout)
for ph in list(slide.placeholders):  # 清掉版式占位符，全部自绘
    ph._element.getparent().remove(ph._element)

def add_text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, space_after=None):
    """runs: [(text, size_pt, bold, color), ...] 每项一段"""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space_after is not None:
            p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb

add_text(slide, 0.44, 0.29, 4.38, 0.29, [("STUDYMATE AGENT", 12, True, C_MUTED)])
add_text(slide, 0.44, 0.71, 12.29, 0.77, [("「个人 + AI 员工」：OPC 理念的产品化样本", 36, True, C_TEXT)])
add_text(slide, 0.44, 1.62, 12.29, 0.40,
         [("星火燎原，单人成军 —— StudyMate 让每位备考者拥有一支全天候在岗的 AI 教研团队", 15, False, C_MUTED)])

def add_card(s, x, title, bullets):
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.15), Inches(5.9), Inches(3.75))
    card.adjustments[0] = 0.035
    card.fill.solid()
    card.fill.fore_color.rgb = C_CARD
    card.line.color.rgb = C_CARD
    card.line.width = Pt(0)
    card.shadow.inherit = False
    add_text(s, x + 0.30, 2.45, 5.30, 0.44, [(title, 20.25, True, C_TEXT)])
    add_text(s, x + 0.30, 3.15, 5.32, 2.55,
             [(t, 15, False, C_TEXT) for t in bullets], space_after=10)

add_card(slide, 0.44, "传统备考服务", [
    "教研、出题、批改、督学由多角色人力承担",
    "按班级统一节奏推进，难以一人一档",
    "服务成本随人数线性上升，价格居高不下",
])
add_card(slide, 6.59, "StudyMate 的 AI 员工团队", [
    "6 个领域 Agent 在岗：规划、出题、批改、归因、调度、陪伴",
    "资料与掌握度驱动，一人一档、一题一归因",
    "新增用户边际成本 ≈ 模型调用费，趋近于零",
])

add_text(slide, 0.44, 6.18, 12.29, 0.40,
         [("生产侧同样如此：本项目从代码、测试到本计划书，均由「人 + AI 员工」协作流水线产出。", 12.75, False, C_MUTED)])
div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.44), Inches(6.79), Inches(12.46), Inches(0.012))
div.fill.solid(); div.fill.fore_color.rgb = C_DIVIDER
div.line.color.rgb = C_DIVIDER; div.line.width = Pt(0)
div.shadow.inherit = False
add_text(slide, 12.29, 6.98, 0.58, 0.25, [("11", 11.25, False, C_MUTED)])

# ── 4. 移动新页到团队页之前（index 10）─────────────────────────
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
sldIdLst.remove(ids[-1])
sldIdLst.insert(10, ids[-1])

# ── 5. 页码顺延：团队页 11→12，结尾页 12→13 ────────────────────
def renumber(slide_obj, old, new):
    for sh in slide_obj.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == old:
            sh.text_frame.paragraphs[0].runs[0].text = new

renumber(prs.slides[11], "11", "12")  # 团队页（移动后 index 11）
renumber(prs.slides[12], "12", "13")  # 结尾页（移动后 index 12）

prs.save(DST)
print("saved:", DST, "slides:", len(prs.slides))
